#!/usr/bin/env python3
"""Verify text formatting in PowerPoint slides.

Usage: python verify_formatting.py <pptx_path> [--caption-only]

Checks font name, size, bold, color, and position for all text shapes.
Use --caption-only to filter for likely captions (long text, lower position).
"""

import sys
import argparse
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor


def is_likely_caption(shape, min_length=40, min_top_inches=5):
    """Heuristic to identify content captions vs UI elements."""
    if not hasattr(shape, 'text_frame'):
        return False
    text = shape.text.strip()
    if len(text) < min_length:
        return False
    try:
        top_inches = Emu(shape.top).inches
        return top_inches >= min_top_inches
    except (AttributeError, NotImplementedError):
        return False


def format_shape_info(shape, slide_num):
    """Extract formatted information about a shape."""
    lines = []
    text = shape.text.strip() if hasattr(shape, 'text') else ''
    if not text:
        return None

    preview = text[:60] + '...' if len(text) > 60 else text
    lines.append(f"\nSlide {slide_num}: '{preview}'")

    if hasattr(shape, 'left'):
        try:
            lines.append(f"  Position: left={Emu(shape.left).inches:.2f}in, top={Emu(shape.top).inches:.2f}in")
            lines.append(f"  Size: width={Emu(shape.width).inches:.2f}in, height={Emu(shape.height).inches:.2f}in")
        except (AttributeError, NotImplementedError):
            pass

    if hasattr(shape, 'text_frame'):
        tf = shape.text_frame
        for p_idx, para in enumerate(tf.paragraphs):
            font = para.font
            lines.append(f"  Paragraph {p_idx + 1}:")
            lines.append(f"    Font name: {font.name}")
            lines.append(f"    Font size: {font.size.pt if font.size else 'N/A'} pt")
            lines.append(f"    Bold: {font.bold}")

            if font.color.type is not None and hasattr(font.color, 'rgb') and font.color.rgb:
                lines.append(f"    Font color: #{font.color.rgb}")
            else:
                lines.append(f"    Font color: N/A")

    return '\n'.join(lines)


def main():
    parser = argparse.ArgumentParser(description='Verify PowerPoint formatting')
    parser.add_argument('pptx_path', help='Path to .pptx file')
    parser.add_argument('--caption-only', action='store_true',
                        help='Only show likely captions (long text, lower position)')
    args = parser.parse_args()

    prs = Presentation(args.pptx_path)

    print(f"Verification of: {args.pptx_path}")
    print("=" * 70)

    found_any = False
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if args.caption_only and not is_likely_caption(shape):
                continue

            info = format_shape_info(shape, slide_idx)
            if info:
                print(info)
                found_any = True

    if not found_any:
        print("No matching shapes found.")
    else:
        print("\n" + "=" * 70)
        print("Verification complete!")


if __name__ == '__main__':
    main()