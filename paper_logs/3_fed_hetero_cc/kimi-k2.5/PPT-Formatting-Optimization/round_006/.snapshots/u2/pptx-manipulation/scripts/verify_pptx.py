#!/usr/bin/env python3
"""Verify PPTX text and formatting.

Handles multi-run text, grouped shapes, and safe color access.

Usage:
    python verify_pptx.py <pptx_path> [--caption-only] [--compare <other_pptx>]
"""

import sys
import argparse
from pptx import Presentation
from pptx.util import Emu, Inches
from pptx.dml.color import RGBColor


def get_full_text(shape):
    """Extract full text from shape, concatenating all runs."""
    if not hasattr(shape, 'text_frame'):
        return ''
    texts = []
    for para in shape.text_frame.paragraphs:
        para_text = ''.join(run.text for run in para.runs)
        texts.append(para_text)
    return '\n'.join(texts)


def is_likely_caption(shape, min_length=40, min_top_inches=5):
    """Identify content captions vs UI elements."""
    if not hasattr(shape, 'text_frame'):
        return False
    text = shape.text.strip()
    if len(text) < min_length:
        return False
    try:
        top_inches = Emu(shape.top).inches
        return top_inches >= min_top_inches
    except (AttributeError, NotImplementedError):
        return False


def format_font_info(paragraph, p_idx):
    """Format font properties for a paragraph."""
    font = paragraph.font
    lines = [f"  Paragraph {p_idx + 1}:"]

    # Font name
    lines.append(f"    Font name: {font.name or 'N/A'}")

    # Font size
    if font.size:
        lines.append(f"    Font size: {font.size.pt} pt")
    else:
        lines.append("    Font size: N/A")

    # Bold
    lines.append(f"    Bold: {font.bold}")

    # Color (safe access)
    if font.color.type is not None:
        try:
            lines.append(f"    Font color: #{font.color.rgb}")
        except AttributeError:
            lines.append("    Font color: N/A")
    else:
        lines.append("    Font color: N/A (theme)")

    return '\n'.join(lines)


def format_shape_info(shape, slide_num, show_text=True):
    """Format information about a shape."""
    lines = []

    # Shape name and type (safe access)
    try:
        tag = shape._element.tag.split('}')[-1]
    except:
        tag = 'unknown'
    lines.append(f"\nSlide {slide_num}: '{shape.name}' (tag={tag})")

    # Position
    if hasattr(shape, 'left'):
        try:
            lines.append(f"  Position: left={Emu(shape.left).inches:.2f}in, top={Emu(shape.top).inches:.2f}in")
            lines.append(f"  Size: width={Emu(shape.width).inches:.2f}in, height={Emu(shape.height).inches:.2f}in")
        except (AttributeError, NotImplementedError):
            pass

    # Text and font info
    if hasattr(shape, 'text_frame') and show_text:
        tf = shape.text_frame
        for p_idx, para in enumerate(tf.paragraphs):
            lines.append(format_font_info(para, p_idx))
            full_text = ''.join(run.text for run in para.runs)
            if full_text.strip():
                preview = full_text[:50] + '...' if len(full_text) > 50 else full_text
                lines.append(f"    Text: '{preview}'")

    return '\n'.join(lines)


def list_shapes_recursive(slide, indent=0):
    """List shapes, handling grouped shapes."""
    prefix = '  ' * indent
    for shape in slide.shapes:
        try:
            tag = shape._element.tag.split('}')[-1]
        except:
            tag = 'unknown'
        print(f"{prefix}name='{shape.name}' tag={tag}")
        if hasattr(shape, 'shapes'):  # GroupShape
            list_shapes_recursive(shape.shapes, indent + 1)


def main():
    parser = argparse.ArgumentParser(description='Verify PPTX formatting')
    parser.add_argument('pptx_path', help='Path to .pptx file')
    parser.add_argument('--caption-only', action='store_true',
                        help='Only show likely captions (long text, lower position)')
    parser.add_argument('--list-shapes', action='store_true',
                        help='Just list shape names and types')
    parser.add_argument('--compare', metavar='OTHER_PPTX',
                        help='Compare with another PPTX file')
    args = parser.parse_args()

    prs = Presentation(args.pptx_path)

    print(f"=== {args.pptx_path} ===")
    print(f"Slides: {len(prs.slides)}")

    if args.list_shapes:
        for i, slide in enumerate(prs.slides, 1):
            print(f"\nSlide {i}:")
            list_shapes_recursive(slide)
        return

    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if args.caption_only and not is_likely_caption(shape):
                continue
            print(format_shape_info(shape, slide_idx))

    if args.compare:
        print(f"\n=== Comparing with {args.compare} ===")
        prs2 = Presentation(args.compare)
        for slide_idx, (s1, s2) in enumerate(zip(prs.slides, prs2.slides), 1):
            for shape1, shape2 in zip(s1.shapes, s2.shapes):
                if hasattr(shape1, 'text_frame') and hasattr(shape2, 'text_frame'):
                    t1 = get_full_text(shape1)
                    t2 = get_full_text(shape2)
                    if t1 != t2:
                        print(f"Slide {slide_idx} '{shape1.name}': text differs")


if __name__ == '__main__':
    main()
