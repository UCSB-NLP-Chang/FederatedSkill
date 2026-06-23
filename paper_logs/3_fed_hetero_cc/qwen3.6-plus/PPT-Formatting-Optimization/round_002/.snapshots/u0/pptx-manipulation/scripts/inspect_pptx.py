#!/usr/bin/env python3
"""Safely inspect a .pptx file structure, text, and formatting."""
import sys
from pptx import Presentation

def safe_color(run):
    try:
        return str(run.font.color.rgb)
    except AttributeError:
        return "None"

def inspect(path):
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, 1):
        print(f"\n=== Slide {i} ===")
        title = slide.shapes.title.text if slide.shapes.title else "(No Title)"
        print(f"  Title: {title}")
        for shape in slide.shapes:
            print(f"  Shape: {shape.shape_type}, name='{shape.name}', left={shape.left}, top={shape.top}, w={shape.width}, h={shape.height}")
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        print(f"    Run: text='{run.text[:60]}...', font={run.font.name}, size={run.font.size}, color={safe_color(run)}, bold={run.font.bold}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: inspect_pptx.py <path_to_pptx>")
        sys.exit(1)
    inspect(sys.argv[1])
