#!/usr/bin/env python3
"""Verify PPTX text and formatting, handling multi-run text and group shapes safely.

Usage: python verify_pptx.py <file.pptx> [--caption-only]
"""
import sys
from pptx import Presentation
from pptx.util import Emu, Inches

def get_full_text(shape):
    """Extract full text from a shape, concatenating all runs."""
    texts = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            texts.append(run.text)
    return ''.join(texts)

def is_caption(shape, min_len=40, min_top=5):
    """Identify captions: long text in lower portion of slide."""
    if not hasattr(shape, 'text_frame'):
        return False
    text = get_full_text(shape).strip()
    if len(text) < min_len:
        return False
    try:
        return Emu(shape.top).inches >= min_top
    except (AttributeError, NotImplementedError):
        return False

def list_shapes(slide, indent=0):
    """Recursively list shapes, handling groups safely."""
    for shape in slide.shapes:
        prefix = "  " * indent
        tag = shape._element.tag.split("}")[-1]
        print(f"{prefix}name='{shape.name}' tag={tag}")
        if hasattr(shape, "shapes"):
            for child in shape.shapes:
                child_tag = child._element.tag.split("}")[-1]
                print(f"{prefix}  child: '{child.name}' tag={child_tag}")

def show_font_info(shape, slide_num):
    """Display font formatting for a shape."""
    text = get_full_text(shape).strip()
    if not text:
        return
    preview = text[:60] + '...' if len(text) > 60 else text
    print(f"\nSlide {slide_num}: '{preview}'")

    tf = shape.text_frame
    for p_idx, para in enumerate(tf.paragraphs):
        font = para.font
        print(f"  Paragraph {p_idx + 1}:")
        print(f"    Font: {font.name}, {font.size.pt if font.size else 'N/A'}pt")
        print(f"    Bold: {font.bold}")
        if font.color.type is not None:
            print(f"    Color: #{font.color.rgb}")
        else:
            print(f"    Color: N/A")

def main():
    if len(sys.argv) < 2:
        print("Usage: verify_pptx.py <file.pptx> [--caption-only]")
        sys.exit(1)

    path = sys.argv[1]
    caption_only = '--caption-only' in sys.argv

    print(f"\n=== {path} ===")
    prs = Presentation(path)

    for i, slide in enumerate(prs.slides, 1):
        print(f"Slide {i}:")
        list_shapes(slide)

        for shape in slide.shapes:
            if caption_only and not is_caption(shape):
                continue
            if hasattr(shape, 'text_frame'):
                show_font_info(shape, i)

if __name__ == "__main__":
    main()
