#!/usr/bin/env python3
"""Safely inspect a .pptx file structure, handling common python-pptx quirks."""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt

def inspect(path):
    prs = Presentation(path)
    print(f"Slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides, 1):
        print(f"\n=== Slide {i} ===")
        for shape in slide.shapes:
            try:
                stype = shape.shape_type
            except NotImplementedError:
                stype = "UNKNOWN/Group"
            print(f"  Shape: type={stype}, name={shape.name}")
            print(f"    pos: left={shape.left}, top={shape.top}, width={shape.width}, height={shape.height}")
            
            # Handle GroupShape children
            if hasattr(shape, 'shapes'):
                for child in shape.shapes:
                    print(f"    CHILD: type={child.shape_type}, name={child.name}")
                    if child.has_text_frame:
                        for j, para in enumerate(child.text_frame.paragraphs):
                            text = para.text.strip()
                            print(f"      C-Para {j}: text=\"{text}\"")
            elif shape.has_text_frame:
                for j, para in enumerate(shape.text_frame.paragraphs):
                    align = para.alignment
                    text = para.text.strip()
                    print(f"    Para {j}: align={align}, text=\"{text}\"")
                    for run in para.runs:
                        font = run.font
                        color_val = "none"
                        if font.color.type is not None:
                            try:
                                color_val = font.color.rgb
                            except AttributeError:
                                color_val = "unknown_type"
                        print(f"      Run: name={font.name}, size={font.size}, color={color_val}, bold={font.bold}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python inspect_pptx.py <file.pptx>")
        sys.exit(1)
    inspect(sys.argv[1])
