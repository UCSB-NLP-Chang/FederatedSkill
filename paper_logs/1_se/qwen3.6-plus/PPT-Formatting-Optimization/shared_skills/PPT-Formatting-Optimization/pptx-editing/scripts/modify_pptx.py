#!/usr/bin/env python3
"""Scaffold for safely modifying .pptx files with python-pptx.

Usage: python modify_pptx.py <input.pptx> <output.pptx>
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def safe_set_font(run, name=None, size_pt=None, color_hex=None, bold=None, italic=None):
    """Safely update run font properties without triggering AttributeError on _NoneColor."""
    if name is not None:
        run.font.name = name
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color_hex is not None:
        # python-pptx allows direct assignment; it handles color object creation safely
        run.font.color.rgb = RGBColor.from_string(color_hex)

def safe_set_position(shape, left_in=None, top_in=None, width_in=None, height_in=None):
    """Safely update shape position/size using Inches."""
    if left_in is not None:
        shape.left = Inches(left_in)
    if top_in is not None:
        shape.top = Inches(top_in)
    if width_in is not None:
        shape.width = Inches(width_in)
    if height_in is not None:
        shape.height = Inches(height_in)

def main(pptx_path, output_path):
    prs = Presentation(pptx_path)
    
    # Example modification loop: replace with task-specific logic
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        safe_set_font(run, name="Calibri", size_pt=17, color_hex="4A6A54", bold=False)
                        # Example: disable word wrap to force single line
                        shape.text_frame.word_wrap = False
                        shape.text_frame.auto_size = None
    
    prs.save(output_path)
    print(f"Saved to {output_path}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python modify_pptx.py <input.pptx> <output.pptx>")
        sys.exit(1)
    main(sys.argv[1], sys.argv[2])
